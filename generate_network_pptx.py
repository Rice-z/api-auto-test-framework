# ============================================
# 黑龙江大学1号教学楼网络规划项目分析 — PPTX 生成脚本
# ============================================

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ======================== 全局配色 ========================
BLUE_DARK  = RGBColor(0x0D, 0x47, 0xA1)   # 深蓝（标题）
BLUE_MID   = RGBColor(0x15, 0x65, 0xC0)   # 中蓝
BLUE_LIGHT = RGBColor(0x42, 0xA5, 0xF5)   # 浅蓝
BLUE_BG    = RGBColor(0xE3, 0xF2, 0xFD)   # 极浅蓝（背景）
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x33, 0x33, 0x33)
GRAY       = RGBColor(0x75, 0x75, 0x75)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
ORANGE     = RGBColor(0xFF, 0x98, 0x00)
RED        = RGBColor(0xF4, 0x43, 0x36)
YELLOW     = RGBColor(0xFF, 0xEB, 0x3B)

W = Inches(13.333)  # 宽屏 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ======================== 工具函数 ========================

def add_bg(slide, color=WHITE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, text="", font_size=10, font_color=BLACK, bold=False):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_title_bar(slide, title_text, subtitle_text=""):
    """添加页面顶部蓝色标题栏"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), BLUE_DARK)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6), title_text,
                font_size=28, font_color=WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4), subtitle_text,
                    font_size=13, font_color=RGBColor(0xBB, 0xDE, 0xFB))
    # 底部装饰线
    add_rect(slide, Inches(0), Inches(1.15), W, Inches(0.04), BLUE_LIGHT)

def add_page_number(slide, num):
    """右下角页码"""
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
                str(num), font_size=10, font_color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_sub_title(slide, text, top=Inches(1.4)):
    """添加小节标题"""
    add_textbox(slide, Inches(0.6), top, Inches(10), Inches(0.5),
                text, font_size=20, font_color=BLUE_DARK, bold=True)
    add_rect(slide, Inches(0.6), top + Inches(0.45), Inches(1.5), Inches(0.04), BLUE_LIGHT)

def add_section_header(slide, num, title, subtitle=""):
    """章节封面"""
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, BLUE_DARK)
    add_rect(slide, Inches(0.15), Inches(0), W - Inches(0.15), H, BLUE_BG)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(1.2),
                f"{num:02d}", font_size=72, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(3.5), Inches(2.0), Inches(8), Inches(1.0),
                title, font_size=36, font_color=BLUE_DARK, bold=True)
    if subtitle:
        add_textbox(slide, Inches(3.5), Inches(2.9), Inches(8), Inches(0.6),
                    subtitle, font_size=16, font_color=GRAY)

def add_info_card(slide, left, top, width, height, label, value):
    """信息卡片"""
    add_rect(slide, left, top, width, height, WHITE, border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_textbox(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.3),
                label, font_size=10, font_color=GRAY)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.3), width - Inches(0.3), Inches(0.4),
                value, font_size=13, font_color=BLACK, bold=True)

def add_floor_room(slide, left, top, width, height, name, color, text_color=BLACK, font_size=8):
    """绘制平面图房间"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x90, 0x90, 0x90)
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape

def make_table(slide, left, top, col_widths, headers, rows):
    """创建表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          sum(col_widths), Inches(0.35 * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    # 表头
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    # 数据行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
    return table_shape


# ======================== 第1页：封面 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_rect(slide, Inches(0), Inches(0), W, H, BLUE_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE_LIGHT)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), BLUE_LIGHT)

# 装饰矩形
add_rect(slide, Inches(2.0), Inches(1.8), Inches(0.08), Inches(2.0), BLUE_LIGHT)
add_textbox(slide, Inches(2.5), Inches(1.8), Inches(9), Inches(1.0),
            "黑龙江大学1号教学楼", font_size=44, font_color=WHITE, bold=True)
add_textbox(slide, Inches(2.5), Inches(2.7), Inches(9), Inches(0.8),
            "网络规划项目分析", font_size=36, font_color=RGBColor(0x90, 0xCA, 0xF9))
add_rect(slide, Inches(2.5), Inches(3.7), Inches(3.0), Inches(0.03), BLUE_LIGHT)

info_texts = [
    "项目类型：校园教学楼网络规划与设计",
    "建筑规模：地上4层 | 总建筑面积约 8,000 m²",
    "覆盖范围：教室、实验室、办公室、走廊、楼梯间",
    "设计内容：综合布线、网络拓扑、VLAN规划、设备选型、语音播报系统",
]
for i, t in enumerate(info_texts):
    add_textbox(slide, Inches(2.5), Inches(4.1 + i * 0.40), Inches(9), Inches(0.35),
                t, font_size=15, font_color=RGBColor(0xE0, 0xE0, 0xE0))

print("[1/22] 封面完成")

# ======================== 第2页：目录 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "目  录", "CONTENTS")

toc_items = [
    ("01", "项目概况", "项目背景、建设目标、建筑规模、人员结构"),
    ("02", "需求分析", "业务需求、网络规模、安全需求、扩展需求"),
    ("03", "网络总体设计", "网络架构设计、三层模型、网络拓扑图"),
    ("04", "楼层平面图", "F1~F4 各层平面布局、设施分布、实际场景图"),
    ("05", "综合布线系统", "工作区/水平/管理/垂直/设备间/建筑群子系统"),
    ("06", "VLAN与IP规划", "VLAN划分策略、IP地址规划、路由配置"),
    ("07", "语音播报系统", "广播系统架构、走廊/楼梯间/教室部署方案"),
    ("08", "设备选型与预算", "核心/汇聚/接入/安全/无线设备选型及预算"),
    ("09", "实施与运维", "实施计划、验收标准、运维管理方案"),
]

for i, (num, title, desc) in enumerate(toc_items):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.6 + row * 1.75)
    add_rect(slide, x, y, Inches(3.6), Inches(1.45), GRAY_LIGHT if (row + col) % 2 == 0 else WHITE,
             border_color=RGBColor(0xE8, 0xE8, 0xE8))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.5),
                num, font_size=28, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.2), Inches(2.4), Inches(0.4),
                title, font_size=16, font_color=BLACK, bold=True)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.65), Inches(2.4), Inches(0.5),
                desc, font_size=10, font_color=GRAY)

add_page_number(slide, 2)
print("[2/22] 目录完成")

# ======================== 第3页：项目概况（章节封面）=======================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "项目概况", "Project Overview")
add_page_number(slide, 3)
print("[3/22] 项目概况封面")

# ======================== 第4页：项目基本信息 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "项目概况", "黑龙江大学1号教学楼网络规划项目")
add_sub_title(slide, "项目基本信息")

cards = [
    ("项目名称", "黑龙江大学1号教学楼\n网络规划与设计项目"),
    ("项目地址", "黑龙江省哈尔滨市南岗区\n学府路74号黑龙江大学校内"),
    ("建筑规模", "地上4层\n总建筑面积约 8,000 m²"),
    ("建设单位", "黑龙江大学\n网络与信息中心"),
    ("建设周期", "预计 6 个月\n(设计2个月 + 施工4个月)"),
    ("功能定位", "综合教学楼\n教室 + 实验室 + 办公室"),
]
for i, (label, value) in enumerate(cards):
    col = i % 3
    row = i // 3
    add_info_card(slide, Inches(0.6 + col * 4.1), Inches(2.2 + row * 2.1),
                  Inches(3.7), Inches(1.7), label, value)

add_page_number(slide, 4)
print("[4/22] 项目基本信息")

# ======================== 第5页：建筑规模与功能分区 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "项目概况", "建筑规模与功能分区")
add_sub_title(slide, "各楼层功能分布")

# 楼层功能表
headers = ["楼层", "面积(m²)", "主要功能", "信息点数量", "备注"]
rows = [
    ["F1（一层）", "2,000", "大厅、网络中心、阶梯教室、教师办公室", "120", "主配线间所在楼层"],
    ["F2（二层）", "2,000", "标准教室、计算机实验室、会议室", "160", "含 2 间计算机房"],
    ["F3（三层）", "2,000", "标准教室、语言实验室、学术报告厅", "150", "含语音播报控制室"],
    ["F4（四层）", "2,000", "标准教室、科研办公室、小型研讨室", "130", ""],
]
make_table(slide, Inches(0.6), Inches(2.3),
           [Inches(1.5), Inches(1.3), Inches(4.0), Inches(1.3), Inches(3.5)],
           headers, rows)

# 人员规模
add_sub_title(slide, "人员与终端规模", top=Inches(4.5))
person_rows = [
    ["教师 / 行政人员", "80", "每人 1 台有线 + 1 台无线终端"],
    ["学生（同时在线）", "600", "每人 1-2 台无线终端"],
    ["其他设备", "—", "语音广播终端、监控摄像头、打印机等"],
    ["合计终端数", "约 1,500+", "有线约 300 + 无线约 1,200"],
]
make_table(slide, Inches(0.6), Inches(5.0),
           [Inches(3.0), Inches(1.5), Inches(6.5)],
           ["用户类别", "数量", "终端说明"], person_rows)

add_page_number(slide, 5)
print("[5/22] 建筑规模")

# ======================== 第6页：建设目标 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "项目概况", "建设目标与设计原则")
add_sub_title(slide, "建设目标")

goals = [
    ("🎓", "教学支撑", "为教室、实验室提供稳定高速网络，支撑多媒体教学、在线考试、远程授课"),
    ("🔗", "全覆盖接入", "实现有线+无线全楼覆盖，千兆到桌面，万兆主干，无线 Wi-Fi 6 覆盖"),
    ("📢", "语音播报", "建设楼内语音广播系统，覆盖教室、走廊、楼梯间，支持分区广播与应急广播"),
    ("🔒", "安全可控", "出口防火墙 + 行为管理 + VLAN 隔离，保障教学网络与办公网络安全"),
    ("📈", "易扩展", "采用模块化设计，预留 30% 端口余量，支持未来 5 年业务扩展"),
]
for i, (icon, title, desc) in enumerate(goals):
    y = Inches(2.2 + i * 1.0)
    add_rect(slide, Inches(0.6), y, Inches(12.0), Inches(0.85), GRAY_LIGHT if i % 2 == 0 else WHITE,
             border_color=RGBColor(0xE8, 0xE8, 0xE8))
    add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.5), Inches(0.5),
                icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), y + Inches(0.08), Inches(2.0), Inches(0.35),
                title, font_size=15, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(1.4), y + Inches(0.42), Inches(10.5), Inches(0.35),
                desc, font_size=11, font_color=BLACK)

add_page_number(slide, 6)
print("[6/22] 建设目标")

# ======================== 第7页：需求分析（章节封面）=======================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "需求分析", "Requirements Analysis")
add_page_number(slide, 7)
print("[7/22] 需求分析封面")

# ======================== 第8页：业务与网络需求 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "需求分析", "业务需求与网络规模")

add_sub_title(slide, "业务需求")
biz_req = [
    ("多媒体教学", "每间教室需接入投影仪、电子白板、教师PC，带宽≥100Mbps"),
    ("在线考试", "计算机实验室需稳定网络，考试期间禁止外网访问，需 VLAN 隔离"),
    ("远程授课", "学术报告厅支持视频会议，上行带宽≥50Mbps，延迟＜50ms"),
    ("语音广播", "覆盖全楼走廊+教室+楼梯间，支持定时打铃、分区通知、应急广播"),
    ("办公网络", "教师办公室接入校园网，支持 OA、邮件、文件共享"),
    ("无线覆盖", "全楼 Wi-Fi 6 覆盖，走廊、教室、办公室无死角，支持高密接入"),
]
for i, (title, desc) in enumerate(biz_req):
    y = Inches(2.2 + i * 0.62)
    add_textbox(slide, Inches(0.8), y, Inches(3.2), Inches(0.3),
                f"▸ {title}", font_size=12, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(4.2), y, Inches(8.0), Inches(0.3),
                desc, font_size=11, font_color=BLACK)

add_sub_title(slide, "网络规模需求", top=Inches(5.95))
net_scale = [
    ["有线信息点", "560+", "教室、办公室、实验室（含预留 30%）"],
    ["无线 AP 点", "40+", "Wi-Fi 6 AP，走廊+大房间全覆盖"],
    ["语音广播终端", "50+", "IP 网络音箱 + 功放 + 喇叭"],
    ["监控摄像头", "20+", "走廊+出入口+楼梯间"],
]
make_table(slide, Inches(0.6), Inches(6.3),
           [Inches(2.5), Inches(1.5), Inches(6.0)], ["类别", "数量", "说明"], net_scale)

add_page_number(slide, 8)
print("[8/22] 需求分析")

# ======================== 第9页：安全与扩展需求 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "需求分析", "安全需求与扩展需求")
add_sub_title(slide, "安全需求")

sec_items = [
    ("边界安全", "出口部署防火墙 + 入侵检测(IDS)，阻止外部攻击"),
    ("访问控制", "基于 802.1X + MAC 认证，教学网与办公网 VLAN 隔离"),
    ("行为审计", "部署上网行为管理，记录访问日志，满足合规要求"),
    ("物理安全", "网络机柜加锁，UPS 供电≥30min，防雷接地保护"),
]
for i, (title, desc) in enumerate(sec_items):
    y = Inches(2.2 + i * 0.7)
    add_rect(slide, Inches(0.6), y, Inches(0.06), Inches(0.5), BLUE_LIGHT)
    add_textbox(slide, Inches(0.9), y + Inches(0.02), Inches(2.5), Inches(0.3),
                title, font_size=14, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(0.9), y + Inches(0.32), Inches(10), Inches(0.3),
                desc, font_size=11, font_color=GRAY)

add_sub_title(slide, "扩展需求", top=Inches(5.0))
ext_items = [
    ("端口预留", "每台交换机预留 30% 端口，便于新增设备"),
    ("带宽扩展", "主干支持 40Gbps 升级，核心交换机支持堆叠"),
    ("无线扩展", "AP 位置预留网线，支持未来增加 AP 密度"),
    ("IPv6 兼容", "所有设备支持 IPv4/IPv6 双栈，平滑过渡"),
]
for i, (title, desc) in enumerate(ext_items):
    y = Inches(5.5 + i * 0.5)
    add_textbox(slide, Inches(0.9), y, Inches(3.0), Inches(0.25),
                f"● {title}", font_size=12, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(4.0), y, Inches(8.0), Inches(0.25),
                desc, font_size=11, font_color=BLACK)

add_page_number(slide, 9)
print("[9/22] 安全与扩展需求")

# ======================== 第10页：网络总体设计（章节封面）=======================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "网络总体设计", "Network Architecture Design")
add_page_number(slide, 10)
print("[10/22] 网络总体设计封面")

# ======================== 第11页：三层网络架构 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "网络总体设计", "三层网络架构设计")
add_sub_title(slide, "核心层 → 汇聚层 → 接入层")

# 绘制三层架构示意图
core_y = Inches(2.0)
agg_y = Inches(3.3)
acc_y = Inches(4.8)
end_y = Inches(6.0)

# 核心层
add_rect(slide, Inches(4.5), core_y, Inches(4.0), Inches(0.85), BLUE_DARK, text="核心层：2台核心交换机（堆叠）\n万兆互联 | 冗余备份", font_size=12, font_color=WHITE, bold=True)
add_textbox(slide, Inches(0.5), core_y + Inches(0.2), Inches(3.5), Inches(0.4),
            "Core Layer", font_size=13, font_color=BLUE_DARK, bold=True)

# 汇聚层
add_rect(slide, Inches(1.0), agg_y, Inches(2.5), Inches(0.85), BLUE_MID, text="F1 汇聚交换机\n上行万兆", font_size=11, font_color=WHITE, bold=True)
add_rect(slide, Inches(4.0), agg_y, Inches(2.5), Inches(0.85), BLUE_MID, text="F2 汇聚交换机\n上行万兆", font_size=11, font_color=WHITE, bold=True)
add_rect(slide, Inches(7.0), agg_y, Inches(2.5), Inches(0.85), BLUE_MID, text="F3 汇聚交换机\n上行万兆", font_size=11, font_color=WHITE, bold=True)
add_rect(slide, Inches(10.0), agg_y, Inches(2.5), Inches(0.85), BLUE_MID, text="F4 汇聚交换机\n上行万兆", font_size=11, font_color=WHITE, bold=True)
add_textbox(slide, Inches(0.5), agg_y + Inches(0.2), Inches(3.5), Inches(0.4),
            "Distribution Layer", font_size=13, font_color=BLUE_MID, bold=True)

# 接入层
for fi in range(4):
    add_rect(slide, Inches(1.0 + fi * 3.0), acc_y, Inches(2.5), Inches(1.1), BLUE_LIGHT,
             text=f"F{fi+1} 接入交换机 ×3\n48口千兆 | PoE供电\n连接教室/办公室/AP", font_size=9, font_color=WHITE, bold=True)
add_textbox(slide, Inches(0.5), acc_y + Inches(0.3), Inches(3.5), Inches(0.4),
            "Access Layer", font_size=13, font_color=BLUE_LIGHT, bold=True)

# 终端层
add_textbox(slide, Inches(0.5), end_y + Inches(0.1), Inches(3.5), Inches(0.4),
            "Endpoints", font_size=13, font_color=GRAY, bold=True)
endpoints = ["教室PC/投影仪", "Wi-Fi 6 AP", "语音广播终端", "监控摄像头", "打印机/其他"]
for ei, ep in enumerate(endpoints):
    add_rect(slide, Inches(0.8 + ei * 2.5), end_y, Inches(2.2), Inches(0.5), GRAY_LIGHT,
             border_color=RGBColor(0xD0, 0xD0, 0xD0), text=ep, font_size=9, font_color=BLACK)

add_page_number(slide, 11)
print("[11/22] 三层架构")

# ======================== 第12页：网络拓扑图 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "网络总体设计", "网络拓扑图")
add_sub_title(slide, "黑龙江大学1号教学楼 — 全网拓扑结构")

# 绘制拓扑图
# 校园网出口
add_rect(slide, Inches(4.8), Inches(2.0), Inches(3.4), Inches(0.7), RGBColor(0xFF, 0x57, 0x22),
         text="校园网核心 / 互联网出口", font_size=11, font_color=WHITE, bold=True)

# 防火墙
add_rect(slide, Inches(5.3), Inches(2.95), Inches(2.4), Inches(0.55), RED,
         text="防火墙 + IDS", font_size=10, font_color=WHITE, bold=True)

# 核心交换机
add_rect(slide, Inches(3.5), Inches(3.8), Inches(6.0), Inches(0.8), BLUE_DARK,
         text="核心交换机 ×2（堆叠）  |  万兆光纤互联", font_size=12, font_color=WHITE, bold=True)

# 汇聚交换机
for fi in range(4):
    add_rect(slide, Inches(1.2 + fi * 2.85), Inches(5.0), Inches(2.4), Inches(0.65), BLUE_MID,
             text=f"F{fi+1} 汇聚交换机", font_size=10, font_color=WHITE, bold=True)

# 接入交换机
for fi in range(4):
    add_rect(slide, Inches(1.2 + fi * 2.85), Inches(6.0), Inches(2.4), Inches(0.55), BLUE_LIGHT,
             text=f"F{fi+1} 接入交换机 ×3 (PoE)", font_size=9, font_color=WHITE, bold=True)

# 防火墙到核心的连线（用矩形模拟）
add_rect(slide, Inches(5.9), Inches(3.5), Inches(0.04), Inches(0.3), RED)

# 核心到汇聚的连线
for fi in range(4):
    add_rect(slide, Inches(2.4 + fi * 2.85), Inches(4.6), Inches(0.04), Inches(0.4), BLUE_DARK)
    add_rect(slide, Inches(2.4 + fi * 2.85), Inches(5.65), Inches(0.04), Inches(0.35), BLUE_MID)

add_page_number(slide, 12)
print("[12/22] 网络拓扑图")

# ======================== 第13页：楼层平面图（章节封面）=======================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "楼层平面图", "Floor Plans — F1 to F4")
add_page_number(slide, 13)
print("[13/22] 楼层平面图封面")

# ======================== 第14-17页：F1~F4 平面图 ========================
def draw_floor_plan(slide, floor_num, floor_name, rooms, special_rooms=None):
    """绘制楼层平面图"""
    add_bg(slide, WHITE)
    add_title_bar(slide, f"F{floor_num} 平面图 — {floor_name}",
                  f"信息点：有线+无线+语音+监控 | 面积：约2,000m²")
    add_sub_title(slide, "平面布局")

    # 建筑外框
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.2), WHITE,
             border_color=RGBColor(0x40, 0x40, 0x40))

    # 横向走廊 (宽0.6英寸)
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(9.0), Inches(0.55), RGBColor(0xEC, 0xEF, 0xF1),
             text="走  廊（吊顶内走线槽 + 照明 + 广播喇叭 + AP）", font_size=8, font_color=GRAY)

    # 竖向走廊
    add_rect(slide, Inches(4.7), Inches(2.0), Inches(0.6), Inches(5.2), RGBColor(0xEC, 0xEF, 0xF1),
             text="走廊", font_size=7, font_color=GRAY)

    # 楼梯间 (左右各一个)
    # 左楼梯间
    stair_l = add_floor_room(slide, Inches(0.5), Inches(2.0), Inches(1.5), Inches(2.3),
                 f"楼梯间A\n(1→4F)\n\n广播喇叭\n监控摄像头", RGBColor(0xFF, 0xF3, 0xE0), font_size=7)
    add_textbox(slide, Inches(0.55), Inches(3.9), Inches(1.4), Inches(0.3),
                "🚨 消防通道", font_size=7, font_color=RED, bold=True)

    # 右楼梯间
    stair_r = add_floor_room(slide, Inches(8.0), Inches(2.0), Inches(1.5), Inches(2.3),
                 f"楼梯间B\n(1→4F)\n\n广播喇叭\n监控摄像头", RGBColor(0xFF, 0xF3, 0xE0), font_size=7)
    add_textbox(slide, Inches(8.05), Inches(3.9), Inches(1.4), Inches(0.3),
                "🚨 消防通道", font_size=7, font_color=RED, bold=True)

    # 绘制房间
    for room in rooms:
        name, x, y, w, h, color = room
        add_floor_room(slide, Inches(x), Inches(y), Inches(w), Inches(h), name, color, font_size=7)

    # 右侧图例
    legend_x = Inches(9.8)
    legend_items = [
        ("网络中心/主配线间", BLUE_DARK, WHITE),
        ("楼层机房", RGBColor(0xB0, 0xBE, 0xC5), WHITE),
        ("教室/实验室", RGBColor(0xE8, 0xF5, 0xE9), BLACK),
        ("教师办公室", RGBColor(0xE3, 0xF2, 0xFD), BLACK),
        ("楼梯间", RGBColor(0xFF, 0xF3, 0xE0), BLACK),
        ("走廊", RGBColor(0xEC, 0xEF, 0xF1), BLACK),
        ("卫生间", RGBColor(0xFC, 0xE4, 0xEC), BLACK),
    ]
    add_textbox(slide, legend_x, Inches(2.0), Inches(3.0), Inches(0.3),
                "图例", font_size=12, font_color=BLUE_DARK, bold=True)
    for li, (name, color, tc) in enumerate(legend_items):
        add_floor_room(slide, legend_x, Inches(2.4 + li * 0.42), Inches(0.5), Inches(0.35),
                       "", color, font_size=6)
        add_textbox(slide, legend_x + Inches(0.6), Inches(2.42 + li * 0.42), Inches(2.2), Inches(0.3),
                    name, font_size=10, font_color=BLACK)

    # 设施标注
    add_textbox(slide, legend_x, Inches(5.3), Inches(3.0), Inches(0.3),
                f"F{floor_num} 层设施清单", font_size=12, font_color=BLUE_DARK, bold=True)
    facilities = [
        f"• 接入交换机：3台（48口 PoE）",
        f"• Wi-Fi 6 AP：10个（走廊×4 + 房间×6）",
        f"• 语音广播终端：12个",
        f"• 监控摄像头：5个",
        f"• 信息面板（双口）：140-160个",
    ]
    for fi, fac in enumerate(facilities):
        add_textbox(slide, legend_x, Inches(5.7 + fi * 0.28), Inches(3.0), Inches(0.25),
                    fac, font_size=9, font_color=BLACK)

    add_page_number(slide, 13 + floor_num)

# F1 平面图
f1_rooms = [
    ("大厅/前台", 2.2, 2.0, 2.5, 2.3, RGBColor(0xE8, 0xE0, 0xF7)),
    ("网络中心\n(主配线间)\n核心交换机\n防火墙", 2.2, 4.85, 2.5, 2.35, BLUE_DARK),
    ("阶梯教室\n101（200人）", 5.5, 2.0, 2.5, 2.3, RGBColor(0xE8, 0xF5, 0xE9)),
    ("教师办公室\n102-103", 5.5, 4.85, 2.5, 2.35, RGBColor(0xE3, 0xF2, 0xFD)),
    ("卫生间\n(男/女)", 0.5, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
    ("卫生间\n(男/女)", 8.0, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
draw_floor_plan(slide, 1, "大厅·网络中心·阶梯教室·办公室", f1_rooms)
print("[14/22] F1平面图")

# F2 平面图
f2_rooms = [
    ("标准教室 201", 2.2, 2.0, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("标准教室 202", 2.2, 3.2, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("计算机实验室\n203（50台PC）", 2.2, 4.85, 2.5, 2.35, RGBColor(0xC8, 0xE6, 0xC9)),
    ("会议室 205", 5.5, 6.05, 2.5, 1.15, RGBColor(0xE3, 0xF2, 0xFD)),
    # ---- 右侧机房 ----
    ("机房 2A\n（服务器/网络机柜）", 5.5, 2.0, 2.5, 2.3, RGBColor(0xB0, 0xBE, 0xC5)),
    ("机房 2B\n（教学服务器集群）", 5.5, 4.85, 2.5, 1.1, RGBColor(0xB0, 0xBE, 0xC5)),
    ("卫生间", 0.5, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
    ("卫生间", 8.0, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
draw_floor_plan(slide, 2, "标准教室·计算机实验室·会议室·机房", f2_rooms)
print("[15/22] F2平面图")

# F3 平面图
f3_rooms = [
    ("标准教室 301", 2.2, 2.0, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("标准教室 302", 2.2, 3.2, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("语言实验室\n304（语音控制室）", 2.2, 4.85, 2.5, 2.35, RGBColor(0xB3, 0xE5, 0xFC)),
    ("广播控制室\n305（语音播报中心）", 5.5, 6.05, 2.5, 1.15, RGBColor(0xB2, 0xDF, 0xDB)),
    # ---- 右侧机房 ----
    ("机房 3A\n（学术报告厅设备间）", 5.5, 2.0, 1.2, 2.3, RGBColor(0xB0, 0xBE, 0xC5)),
    ("机房 3B\n（语言实验室服务器）", 6.8, 2.0, 1.2, 2.3, RGBColor(0xB0, 0xBE, 0xC5)),
    ("机房 3C\n（广播系统设备柜）", 5.5, 4.85, 2.5, 1.1, RGBColor(0xB0, 0xBE, 0xC5)),
    ("卫生间", 0.5, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
    ("卫生间", 8.0, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
draw_floor_plan(slide, 3, "标准教室·语音实验室·广播控制·机房", f3_rooms)
print("[16/22] F3平面图")

# F4 平面图
f4_rooms = [
    ("标准教室 401", 2.2, 2.0, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("标准教室 402", 2.2, 3.2, 2.5, 1.1, RGBColor(0xE8, 0xF5, 0xE9)),
    ("小型研讨室\n405", 2.2, 4.85, 2.5, 1.1, RGBColor(0xE8, 0xE0, 0xF7)),
    ("研究生工作室\n406（40工位）", 2.2, 6.05, 2.5, 1.15, RGBColor(0xE3, 0xF2, 0xFD)),
    # ---- 右侧机房 ----
    ("机房 4A\n（科研计算服务器）", 5.5, 2.0, 1.2, 2.3, RGBColor(0xB0, 0xBE, 0xC5)),
    ("机房 4B\n（数据存储备份）", 6.8, 2.0, 1.2, 2.3, RGBColor(0xB0, 0xBE, 0xC5)),
    ("机房 4C\n（GPU 计算节点）", 5.5, 4.85, 2.5, 2.35, RGBColor(0xB0, 0xBE, 0xC5)),
    ("卫生间", 0.5, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
    ("卫生间", 8.0, 4.85, 1.5, 2.35, RGBColor(0xFC, 0xE4, 0xEC)),
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
draw_floor_plan(slide, 4, "标准教室·科研办公·机房", f4_rooms)
print("[17/22] F4平面图")

# ======================== 第18页：综合布线系统 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "综合布线系统", "Structured Cabling System")
add_sub_title(slide, "六子系统设计")

cabling = [
    ("工作区子系统", "每个信息面板提供双口（数据+语音），教室/办公室标配 2-4 个面板\n面板类型：RJ45 六类非屏蔽模块"),
    ("水平子系统", "六类非屏蔽双绞线（Cat6 UTP），PoE 供电传输距离 ≤90m\n走廊吊顶内金属线槽敷设，入室采用 PVC 线管"),
    ("管理子系统", "每层弱电间设置机柜（42U），安装接入交换机+配线架\n楼层配线间兼做管理间，跳线管理采用电子标签"),
    ("垂直干线子系统", "万兆多模光纤（OM3/OM4）从 F1 主配线间引至各层管理间\n冗余设计：每层 2 根 12 芯光纤，一主一备"),
    ("设备间子系统", "F1 网络中心（主配线间）：核心交换机、防火墙、服务器机柜\nUPS（30kVA，备电≥30min）+ 精密空调"),
    ("建筑群子系统", "通过 24 芯单模光纤接入校园网主干光缆\n与学校网络中心万兆互联"),
]
for i, (title, desc) in enumerate(cabling):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.2 + row * 2.4)
    add_rect(slide, x, y, Inches(3.9), Inches(2.1), GRAY_LIGHT,
             border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, Inches(3.9), Inches(0.45), BLUE_DARK,
             text=title, font_size=13, font_color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.6), Inches(1.4),
                desc, font_size=9, font_color=BLACK)

add_page_number(slide, 18)
print("[18/22] 综合布线系统")

# ======================== 第19页：VLAN与IP规划 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "VLAN与IP规划", "VLAN & IP Address Planning")
add_sub_title(slide, "VLAN 划分策略")

vlan_headers = ["VLAN ID", "VLAN 名称", "网段", "用途", "网关"]
vlan_rows = [
    ["10", "Management", "10.1.10.0/24", "网络设备管理（交换机/AP/防火墙）", "10.1.10.1"],
    ["20", "Teaching", "10.1.20.0/22", "教室多媒体设备（容纳 1022 个 IP）", "10.1.20.1"],
    ["30", "ComputerLab", "10.1.30.0/23", "计算机实验室（容纳 510 个 IP）", "10.1.30.1"],
    ["40", "Office", "10.1.40.0/23", "教师/行政办公（容纳 510 个 IP）", "10.1.40.1"],
    ["50", "Wireless-Staff", "10.1.50.0/22", "教职工无线网络 (Wi-Fi 6)", "10.1.50.1"],
    ["60", "Wireless-Student", "10.1.60.0/21", "学生无线网络（容纳 2046 个 IP）", "10.1.60.1"],
    ["70", "Broadcast", "10.1.70.0/24", "语音播报系统（IP 音箱/功放）", "10.1.70.1"],
    ["80", "Surveillance", "10.1.80.0/24", "视频监控系统", "10.1.80.1"],
    ["90", "Server", "10.1.90.0/24", "楼内服务器（考试/广播/门禁）", "10.1.90.1"],
]
make_table(slide, Inches(0.3), Inches(2.2),
           [Inches(0.8), Inches(1.6), Inches(1.6), Inches(5.2), Inches(1.4)],
           vlan_headers, vlan_rows)

add_sub_title(slide, "路由与安全策略", top=Inches(5.8))
add_textbox(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.8),
            "• 核心交换机启用三层路由（OSPF），VLAN 间默认隔离，按需开放 ACL 策略\n"
            "• 学生无线 VLAN 仅允许访问互联网，禁止访问教学/办公 VLAN\n"
            "• 教学 VLAN 与办公 VLAN 单向互通（办公 → 教学可访问），考试期间通过 ACL 禁止外网\n"
            "• DHCP Snooping + Dynamic ARP Inspection 防止私接路由器和 ARP 欺骗",
            font_size=10, font_color=BLACK)

add_page_number(slide, 19)
print("[19/22] VLAN与IP规划")

# ======================== 第20页：语音播报系统 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "语音播报系统", "Voice Broadcasting System")
add_sub_title(slide, "系统架构与部署方案")

# 系统架构
add_rect(slide, Inches(3.5), Inches(2.0), Inches(6.0), Inches(0.7), BLUE_DARK,
         text="语音播报服务器（F3 广播控制室）+ IP 网络广播主机", font_size=11, font_color=WHITE, bold=True)

add_rect(slide, Inches(3.5), Inches(3.0), Inches(6.0), Inches(0.6), BLUE_MID,
         text="PoE 接入交换机（VLAN 70）→ 每层弱电间", font_size=11, font_color=WHITE, bold=True)

# 终端设备
terminals = [
    ("走廊 IP 音箱", "12 台", "每层 3 台，走廊中段+两端"),
    ("教室 IP 音箱", "32 台", "每间教室 1 台，嵌入式安装"),
    ("楼梯间 IP 音箱", "8 台", "每个楼梯间每层 1 台"),
    ("室外防水音柱", "4 台", "楼门口×2 + 楼后×2"),
    ("无线麦克风", "2 套", "校长讲话/紧急通知"),
    ("消防联动模块", "1 套", "火灾自动切换紧急广播"),
]
for i, (name, qty, desc) in enumerate(terminals):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(4.0 + row * 1.2)
    add_rect(slide, x, y, Inches(3.9), Inches(0.95), GRAY_LIGHT,
             border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.08), Inches(2.0), Inches(0.3),
                name, font_size=12, font_color=BLUE_DARK, bold=True)
    add_textbox(slide, x + Inches(2.2), y + Inches(0.08), Inches(1.5), Inches(0.3),
                qty, font_size=12, font_color=ORANGE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.5), Inches(0.35),
                desc, font_size=9, font_color=GRAY)

add_page_number(slide, 20)
print("[20/22] 语音播报系统")

# ======================== 第21页：设备选型与预算 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "设备选型与预算", "Equipment Selection & Budget")
add_sub_title(slide, "主要网络设备清单")

eq_headers = ["类别", "设备型号（参考）", "数量", "单价(元)", "小计(元)"]
eq_rows = [
    ["核心交换机", "H3C S6520X-30QC-EI（万兆）", "2", "35,000", "70,000"],
    ["汇聚交换机", "H3C S5560X-30C-EI（万兆上行）", "4", "18,000", "72,000"],
    ["接入交换机", "H3C S5130S-52P-PWR-EI（48口PoE）", "12", "8,500", "102,000"],
    ["防火墙", "H3C SecPath F1000-AI-25", "1", "48,000", "48,000"],
    ["无线控制器", "H3C WX3520X（管理64AP）", "1", "22,000", "22,000"],
    ["Wi-Fi 6 AP", "H3C WA6638（吸顶式）", "40", "2,200", "88,000"],
    ["IP广播主机", "ITC T-7700N", "1", "15,000", "15,000"],
    ["IP网络音箱", "ITC T-7707A", "50", "800", "40,000"],
    ["监控摄像头", "海康威视 DS-2CD2T46WD-I3", "20", "650", "13,000"],
    ["服务器", "H3C UniServer R4900 G6", "2", "35,000", "70,000"],
    ["UPS电源", "山特 3C3 HD 30kVA", "1", "45,000", "45,000"],
    ["机柜(42U)", "图腾 G26642", "6", "3,500", "21,000"],
    ["综合布线", "六类线+光纤+桥架+面板+跳线", "1批", "—", "120,000"],
]
eq_rows.append(["", "", "", "合计", "726,000"])
make_table(slide, Inches(0.2), Inches(2.2),
           [Inches(1.6), Inches(5.0), Inches(0.9), Inches(1.5), Inches(1.5)],
           eq_headers, eq_rows)

add_page_number(slide, 21)
print("[21/22] 设备选型与预算")

# ======================== 第22页：实施计划与结束语 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "实施计划与运维管理", "Implementation & O&M Plan")
add_sub_title(slide, "实施计划")

plan_headers = ["阶段", "时间", "工作内容"]
plan_rows = [
    ["第一阶段：需求调研", "第 1-2 周", "现场勘查、需求确认、方案细化"],
    ["第二阶段：深化设计", "第 3-4 周", "施工图设计、设备采购、管线预埋"],
    ["第三阶段：布线施工", "第 5-10 周", "线缆敷设、信息面板安装、光纤熔接、测试验收"],
    ["第四阶段：设备安装", "第 11-14 周", "机柜安装、交换机/AP/广播/监控等设备上架调试"],
    ["第五阶段：系统联调", "第 15-18 周", "全网络联调、VLAN/路由/安全策略配置、语音广播测试"],
    ["第六阶段：试运行", "第 19-22 周", "试运行 1 个月，压力测试、安全扫描、问题整改"],
    ["第七阶段：竣工验收", "第 23-24 周", "文档交付、用户培训、正式移交"],
]
make_table(slide, Inches(0.4), Inches(2.2),
           [Inches(2.5), Inches(1.5), Inches(8.0)], plan_headers, plan_rows)

add_sub_title(slide, "运维管理", top=Inches(5.3))
add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2),
            "• 网络监控：部署 Zabbix 实时监控网络设备状态、链路流量、AP 负载，异常自动告警\n"
            "• 日常巡检：每周巡检核心设备运行状态、UPS 电池、机房温湿度，每月备份配置\n"
            "• 安全管理：定期更新防火墙规则库，每季度安全扫描与渗透测试\n"
            "• 文档管理：建立网络拓扑、设备配置、IP 分配、运维日志等文档体系",
            font_size=11, font_color=BLACK)

add_page_number(slide, 22)
print("[22/22] 实施计划")

# ======================== 保存 ========================
output_path = "c:/Users/ASUS/Desktop/黑龙江大学1号教学楼网络规划项目分析.pptx"
prs.save(output_path)
print(f"\n✅ 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
